from pptx import Presentation
from pptx.util import Inches, Pt
from summarizer import *
from font_size import *


f = open('summary.txt', 'r', encoding='cp1252')
txt = f.read()
f.close()
print(txt)
title1, title2 = get_title()
print(title1)
print(title2)

# Create Presentation
prs = Presentation()

# Set the slide layout to have a modern background
slide_layout = prs.slide_layouts[6]

# Add a slide for each picture path provided
# picture_paths is a variable that contains a list of file paths for the pictures
# that will be added to each slide in the PowerPoint presentation
# for path in picture_paths:
    # Add a new slide
slide = prs.slides.add_slide(slide_layout)

    # Add the picture to the slide
    # pic = slide.shapes.add_picture(path, Inches(1), Inches(1), height=Inches(5.5))

prs.save('example.pptx')