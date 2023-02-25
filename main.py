import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from summarizer import *
from font_size import *
from PIL import Image

if __name__ == '__main__':

    f = open('summary.txt', 'r', encoding='utf-8')
    txt = f.read()
    f.close()
    print(txt)
    title1, title2 = get_title()
    print(title1)
    print(title2)

    # Create Presentation
    prs = Presentation()

    # Set the slide layout to have a modern background
    slide_layout = prs.slide_layouts[6]

    # Add a slide for each picture path provided
    # picture_paths is a variable that contains a list of file paths for the pictures
    # that will be added to each slide in the PowerPoint presentation
    # for path in picture_paths:
    # Add a new slide
    slide = prs.slides.add_slide(slide_layout)

    # Add the picture to the slide
    # pic = slide.shapes.add_picture(path, Inches(1), Inches(1), height=Inches(5.5))

    shapes = slide.shapes
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(4)

    txBox1 = slide.shapes.add_textbox(left, top, width, height)
    tf1 = txBox1.text_frame
    tf1.text = title1

    if len(title1) > 84:
        tf1.paragraphs[0].font.size = Pt(40)
    else:
        tf1.paragraphs[0].font.size = Pt(50)

    top = Inches(5)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.text = title2
    tf2.paragraphs[0].font.size = Pt(20)

    # img_path = 'image.png'
    # pic = slide.shapes.add_picture(img_path, Inches(1), Inches(3), height=Inches(2))
    txBox1.text_frame.word_wrap = True
    txBox2.text_frame.word_wrap = True
    txBox1.text_frame.auto_size = True
    txBox2.text_frame.auto_size = True

    # Background Images
    img_path = "database/picture_2.png"
    im = Image.open('/var/www/examples/heroine.png')
    im.putalpha(100)
    left = top = Inches(0)
    pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

    # This moves it to the background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

    summarizer(prs)

    # for i in range(counter):
    #     new_slide = prs.slides.add_slide(slide_layout)

    prs.save('example.pptx')




