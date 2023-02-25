import collections 
import collections.abc 
from pptx import Presentation
from pptx.util import Inches, Pt
from font_size import *
import requests
from bs4 import BeautifulSoup
from io import BytesIO
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from bing_image_urls import bing_image_urls
from PIL import Image

# nltk.download('averaged_perceptron_tagger')
# nltk.download('punkt')
# creates a new ppt slide, taking the presentation to be added to and text to be added in the slide as parameters
def new_slide(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(4)

    txBox1 = slide.shapes.add_textbox(left, top, width, height)
    tf1 = txBox1.text_frame
    tf1.text = text

    tf1.paragraphs[0].font.size = Pt(24)

    txBox1.text_frame.word_wrap = True
    txBox1.text_frame.auto_size = True

    # Tokenize the sentence into words
    words = word_tokenize(text)

    # Remove stop words from the list of words
    stop_words = set(stopwords.words('english'))
    words = [item for item in words if item.lower() not in stop_words]

    # Get the part of speech for each word
    pos_tags = nltk.pos_tag(words)

    # Extract keywords based on the part of speech
    keywords = [word for word, pos in pos_tags if pos in ['NN', 'NNS', 'NNP', 'NNPS']]
    if keywords == []:
        return
    if "prudence" in keywords:
        keywords.remove("prudence")

    index = 0
    if len(keywords) > 3:
        repeats = 3
    else:
        repeats = len(keywords)
    for i in range(repeats):
        image, url = newPicture(keywords[i])
        if image == None:
            return
        response = requests.get(url)
        img_data = response.content
        left = Inches(0)
        top = Inches(0)
        # try:
        #     pic = slide.shapes.add_picture(image, left, top)
        # except Image.UnidentifiedImageError as e:
        #     print(e)
        #     return
        if url.split(".")[-1] == "jpg":

            try:
                with open(f"database/{keywords[i]}.jpg", "wb") as file:
                    file.write(img_data)
                jpg_image = Image.open(f"database/{keywords[i]}.jpg")
                im = jpg_image.convert("RGBA")
                im.save(f"database/{keywords[i]}.png")
            except (AttributeError, Image.UnidentifiedImageError, IndexError, FileNotFoundError) as e:
                print(e)
                if i > repeats-2:
                    return
                continue
            index = i


    path = f"database/{keywords[index]}.png"
    with open(f"database/{keywords[index]}.png", "wb") as file:
        file.write(img_data)
    # pic.save(path)

    img = Image.open(f"database/{keywords[index]}.png").convert("RGBA")
    mask = Image.new('RGBA', img.size, (255, 255, 255, 0))
    mask.putalpha(150)
    im = Image.composite(img, mask, mask)
    im.save(path)
    left = top = Inches(0)
    pic = slide.shapes.add_picture(path, left, top, width=prs.slide_width, height=prs.slide_height)

    # move to back
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

# searches the internet using "keyword" and returns the picture to be put in the slide
def newPicture(keyword):
    # url = f"https://www.google.com/search?q={keyword}&tbm=isch"
    for i in range(3):
        try:
            url = bing_image_urls(keyword, limit=3)[i]
        except (AttributeError, IndexError) as e:
            print(e)
            if i > 1:
                return [None, None]


    # Send a request to the search URL and get the HTML response
    # headers = {"User-Agent": "Mozilla/5.0"}
    # response = requests.get(url, headers=headers)
    # soup = BeautifulSoup(response.content, "html.parser")

    # Find the first image in the search results and get its URL
    # img_url = soup.find("img")["src"]
    # print(img_url)

    # Download the image from the URL
    try:
        response = requests.get(url)
    except (AttributeError) as e:
        print(e)
        return [None, None]
    img_data = response.content

    # Create an image object using BytesIO
    image = BytesIO(img_data)

    # Return the image object
    return image, url




